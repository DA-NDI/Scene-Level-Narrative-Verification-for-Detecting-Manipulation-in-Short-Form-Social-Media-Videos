from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


INK = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
TEAL = RGBColor(15, 118, 110)
AMBER = RGBColor(180, 83, 9)
BG = RGBColor(248, 250, 252)


def _clean_markdown_line(line: str) -> str:
    txt = line.replace("**", "").replace("*", "").replace("$", "")
    txt = re.sub(r"\s+", " ", txt).strip()
    return txt


def _add_title(slide, title: str, is_cover: bool) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.9), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(50 if is_cover else 40)
    p.font.bold = True
    p.font.color.rgb = INK


def _add_body(slide, body_lines: list[str], is_cover: bool) -> None:
    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.1), Inches(4.7))
    tf = body_box.text_frame
    tf.clear()
    tf.word_wrap = True

    first_written = False
    for raw in body_lines:
        line = raw.strip()
        if not line:
            continue

        level = 0
        if line.startswith("## "):
            txt = _clean_markdown_line(line[3:])
            color = TEAL
            size = Pt(28 if is_cover else 24)
            bold = True
        elif line.startswith("- "):
            txt = _clean_markdown_line(line[2:])
            color = SLATE
            size = Pt(24 if is_cover else 22)
            bold = False
        elif re.match(r"^\d+\.\s+", line):
            txt = _clean_markdown_line(re.sub(r"^\d+\.\s+", "", line))
            color = SLATE
            size = Pt(24 if is_cover else 22)
            bold = False
        else:
            txt = _clean_markdown_line(line)
            if txt.lower().startswith(("problem:", "research gap:", "constraint:", "trade-off:", "interpretation:", "main bottleneck:")):
                color = AMBER
                size = Pt(23 if is_cover else 21)
                bold = True
            else:
                color = SLATE
                size = Pt(23 if is_cover else 21)
                bold = False

        p = tf.paragraphs[0] if not first_written else tf.add_paragraph()
        first_written = True
        p.text = txt
        p.level = level
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.line_spacing = 1.2


def _style_slide_background(slide, prs) -> None:
    bg = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        prs.slide_width,
        prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    # Top accent stripe
    stripe = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.18),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()

    # Thin divider under title area
    divider = slide.shapes.add_shape(
        1,
        Inches(0.7),
        Inches(1.62),
        Inches(11.9),
        Inches(0.03),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(203, 213, 225)
    divider.line.fill.background()

    # Send decorative shapes behind textboxes
    bg.z_order = 0
    stripe.z_order = 1
    divider.z_order = 1


def _add_footer(slide, idx: int) -> None:
    footer = slide.shapes.add_textbox(Inches(11.7), Inches(6.65), Inches(0.8), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(idx)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(100, 116, 139)


def convert(src_path: Path, out_path: Path) -> int:
    text = src_path.read_text(encoding="utf-8")

    if text.startswith("---"):
        parts = text.split("\n---\n", 1)
        if len(parts) == 2:
            text = parts[1]

    slides_raw = [s.strip() for s in re.split(r"\n---\n", text) if s.strip()]

    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000

    for idx, raw in enumerate(slides_raw, start=1):
        lines = [ln.rstrip() for ln in raw.splitlines() if ln.strip()]
        if not lines:
            continue

        title = None
        body = []

        for ln in lines:
            if ln.startswith("# "):
                title = ln[2:].strip()
            elif ln.startswith("## ") and title is None:
                title = ln[3:].strip()
            else:
                body.append(ln)

        if title is None:
            title = "Slide"

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _style_slide_background(slide, prs)
        _add_title(slide, title, is_cover=(idx == 1))
        _add_body(slide, body, is_cover=(idx == 1))
        _add_footer(slide, idx)

    prs.save(str(out_path))
    return len(prs.slides)


if __name__ == "__main__":
    src = Path("master seminar/symposium_2026_slides_starter.marp.md")
    out = Path("master seminar/symposium_2026_slides_styled.pptx")
    count = convert(src, out)
    print(f"WROTE {out}")
    print(f"SLIDES {count}")
